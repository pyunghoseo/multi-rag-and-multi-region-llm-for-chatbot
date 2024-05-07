import json
import boto3
import os
import time
import datetime
from io import BytesIO
import PyPDF2
import csv
import traceback
import re
from urllib import parse

from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.docstore.document import Document
from langchain.chains.summarize import load_summarize_chain
from langchain.chains import ConversationChain
from langchain.memory import ConversationBufferWindowMemory
from langchain.callbacks.streaming_stdout import StreamingStdOutCallbackHandler
from botocore.config import Config

from langchain.vectorstores.faiss import FAISS
from langchain.vectorstores.opensearch_vector_search import OpenSearchVectorSearch
from langchain.embeddings import BedrockEmbeddings
from langchain.retrievers import AmazonKendraRetriever
from multiprocessing import Process, Pipe

from langchain_community.chat_models import BedrockChat
from langchain_core.prompts import MessagesPlaceholder, ChatPromptTemplate

s3 = boto3.client('s3')
s3_bucket = os.environ.get('s3_bucket') # bucket name
s3_prefix = os.environ.get('s3_prefix')
callLogTableName = os.environ.get('callLogTableName')
kendra_region = os.environ.get('kendra_region', 'us-west-2')

profile_of_LLMs = json.loads(os.environ.get('profile_of_LLMs'))
isReady = False   
isDebugging = False

opensearch_account = os.environ.get('opensearch_account')
opensearch_passwd = os.environ.get('opensearch_passwd')
enableReference = os.environ.get('enableReference', 'false')
debugMessageMode = os.environ.get('debugMessageMode', 'false')
opensearch_url = os.environ.get('opensearch_url')
path = os.environ.get('path')
useParallelUpload = os.environ.get('useParallelUpload', 'false')
useParallelRAG = os.environ.get('useParallelRAG', 'false')
kendraIndex = os.environ.get('kendraIndex')
kendra_method = "kendra_retriever" # custom_retriever or kendra_retriever
roleArn = os.environ.get('roleArn')
numberOfRelevantDocs = os.environ.get('numberOfRelevantDocs', '10')
top_k = int(numberOfRelevantDocs)
selected_LLM = 0
capabilities = json.loads(os.environ.get('capabilities'))
print('capabilities: ', capabilities)
MSG_LENGTH = 100

# websocket
connection_url = os.environ.get('connection_url')
client = boto3.client('apigatewaymanagementapi', endpoint_url=connection_url)
print('connection_url: ', connection_url)

HUMAN_PROMPT = "\n\nHuman:"
AI_PROMPT = "\n\nAssistant:"

map_chain = dict() 

def get_chat(profile_of_LLMs, selected_LLM):
    profile = profile_of_LLMs[selected_LLM]
    bedrock_region =  profile['bedrock_region']
    modelId = profile['model_id']
    print(f'LLM: {selected_LLM}, bedrock_region: {bedrock_region}, modelId: {modelId}')
    maxOutputTokens = int(profile['maxOutputTokens'])
                          
    # bedrock   
    boto3_bedrock = boto3.client(
        service_name='bedrock-runtime',
        region_name=bedrock_region,
        config=Config(
            retries = {
                'max_attempts': 30
            }            
        )
    )
    parameters = {
        "max_tokens":maxOutputTokens,     
        "temperature":0.1,
        "top_k":250,
        "top_p":0.9,
        "stop_sequences": [HUMAN_PROMPT]
    }
    # print('parameters: ', parameters)

    chat = BedrockChat(
        model_id=modelId,
        client=boto3_bedrock, 
        streaming=True,
        callbacks=[StreamingStdOutCallbackHandler()],
        model_kwargs=parameters,
    )        
    
    return chat

def get_embedding(profile_of_LLMs, selected_LLM):
    profile = profile_of_LLMs[selected_LLM]
    bedrock_region =  profile['bedrock_region']
    modelId = profile['model_id']
    print(f'Embedding: {selected_LLM}, bedrock_region: {bedrock_region}, modelId: {modelId}')
    
    # bedrock   
    boto3_bedrock = boto3.client(
        service_name='bedrock-runtime',
        region_name=bedrock_region,
        config=Config(
            retries = {
                'max_attempts': 30
            }            
        )
    )
    
    bedrock_embedding = BedrockEmbeddings(
        client=boto3_bedrock,
        region_name = bedrock_region,
        model_id = 'amazon.titan-embed-text-v1' 
    )  
    
    return bedrock_embedding

def sendMessage(id, body):
    try:
        client.post_to_connection(
            ConnectionId=id, 
            Data=json.dumps(body)
        )
    except Exception:
        err_msg = traceback.format_exc()
        print('err_msg: ', err_msg)
        raise Exception ("Not able to send a message")

def sendResultMessage(connectionId, requestId, msg):    
    result = {
        'request_id': requestId,
        'msg': msg,
        'status': 'completed'
    }
    #print('debug: ', json.dumps(debugMsg))
    sendMessage(connectionId, result)
        
def sendDebugMessage(connectionId, requestId, msg):
    debugMsg = {
        'request_id': requestId,
        'msg': msg,
        'status': 'debug'
    }
    #print('debug: ', json.dumps(debugMsg))
    sendMessage(connectionId, debugMsg)

def sendErrorMessage(connectionId, requestId, msg):
    errorMsg = {
        'request_id': requestId,
        'msg': msg,
        'status': 'error'
    }
    print('error: ', json.dumps(errorMsg))
    sendMessage(connectionId, errorMsg)

def general_conversation(connectionId, requestId, chat, query):
    global time_for_inference, history_length, token_counter_history    
    time_for_inference = history_length = token_counter_history = 0
    
    if debugMessageMode == 'true':  
        start_time_for_inference = time.time()
    
    if isKorean(query)==True :
        system = (
            "다음의 Human과 Assistant의 친근한 이전 대화입니다. Assistant은 상황에 맞는 구체적인 세부 정보를 충분히 제공합니다. Assistant의 이름은 서연이고, 모르는 질문을 받으면 솔직히 모른다고 말합니다."
        )
    else: 
        system = (
            "Using the following conversation, answer friendly for the newest question. If you don't know the answer, just say that you don't know, don't try to make up an answer. You will be acting as a thoughtful advisor."
        )
    
    human = "{input}"
    
    prompt = ChatPromptTemplate.from_messages([("system", system), MessagesPlaceholder(variable_name="history"), ("human", human)])
    print('prompt: ', prompt)
    
    history = memory_chain.load_memory_variables({})["chat_history"]
    print('memory_chain: ', history)
                
    chain = prompt | chat    
    try: 
        isTyping(connectionId, requestId)  
        stream = chain.invoke(
            {
                "history": history,
                "input": query,
            }
        )
        msg = readStreamMsg(connectionId, requestId, stream.content)    
                            
        msg = stream.content
        print('msg: ', msg)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)        
            
        sendErrorMessage(connectionId, requestId, err_msg)    
        raise Exception ("Not able to request to LLM")

    if debugMessageMode == 'true':  
        chat_history = ""
        for dialogue_turn in history:
            #print('type: ', dialogue_turn.type)
            #print('content: ', dialogue_turn.content)
            
            dialog = f"{dialogue_turn.type}: {dialogue_turn.content}\n"            
            chat_history = chat_history + dialog
                
        history_length = len(chat_history)
        print('chat_history length: ', history_length)
        
        token_counter_history = 0
        if chat_history:
            token_counter_history = chat.get_num_tokens(chat_history)
            print('token_size of history: ', token_counter_history)
        
        end_time_for_inference = time.time()
        time_for_inference = end_time_for_inference - start_time_for_inference
        
    return msg



def store_document_for_faiss(docs, vectorstore_faiss):
    print('store document into faiss')    
    vectorstore_faiss.add_documents(docs)       
    print('uploaded into faiss')

def store_document_for_opensearch(bedrock_embedding, docs, userId, documentId):
    new_vectorstore = OpenSearchVectorSearch(
        index_name="rag-index-"+userId,
        is_aoss = False,
        #engine="faiss",  # default: nmslib
        embedding_function = bedrock_embedding,
        opensearch_url = opensearch_url,
        http_auth=(opensearch_account, opensearch_passwd),
    )
    response = new_vectorstore.add_documents(docs)    
    print('response of adding documents: ', response)
    
    print('uploaded into opensearch')

# store document into Kendra
def store_document_for_kendra(path, s3_file_name, documentId):
    print('store document to kendra')
    encoded_name = parse.quote(s3_file_name)
    source_uri = path + encoded_name    
    #print('source_uri: ', source_uri)
    ext = (s3_file_name[s3_file_name.rfind('.')+1:len(s3_file_name)]).upper()
    print('ext: ', ext)

    # PLAIN_TEXT, XSLT, MS_WORD, RTF, CSV, JSON, HTML, PDF, PPT, MD, XML, MS_EXCEL
    if(ext == 'PPTX'):
        file_type = 'PPT'
    elif(ext == 'TXT'):
        file_type = 'PLAIN_TEXT'         
    elif(ext == 'XLS' or ext == 'XLSX'):
        file_type = 'MS_EXCEL'      
    elif(ext == 'DOC' or ext == 'DOCX'):
        file_type = 'MS_WORD'
    else:
        file_type = ext

    kendra_client = boto3.client(
        service_name='kendra', 
        region_name=kendra_region,
        config = Config(
            retries=dict(
                max_attempts=10
            )
        )
    )

    documents = [
        {
            "Id": documentId,
            "Title": s3_file_name,
            "S3Path": {
                "Bucket": s3_bucket,
                "Key": s3_prefix+'/'+s3_file_name
            },
            "Attributes": [
                {
                    "Key": '_source_uri',
                    'Value': {
                        'StringValue': source_uri
                    }
                },
                {
                    "Key": '_language_code',
                    'Value': {
                        'StringValue': "ko"
                    }
                },
            ],
            "ContentType": file_type
        }
    ]
    print('document info: ', documents)

    result = kendra_client.batch_put_document(
        IndexId = kendraIndex,
        RoleArn = roleArn,
        Documents = documents       
    )
    # print('batch_put_document(kendra): ', result)
    print('uploaded into kendra')

# load documents from s3 for pdf and txt
def load_document(file_type, s3_file_name):
    s3r = boto3.resource("s3")
    doc = s3r.Object(s3_bucket, s3_prefix+'/'+s3_file_name)
    
    if file_type == 'pdf':
        Byte_contents = doc.get()['Body'].read()
        reader = PyPDF2.PdfReader(BytesIO(Byte_contents))
        
        texts = []
        for page in reader.pages:
            texts.append(page.extract_text())
        contents = '\n'.join(texts)
        
    elif file_type == 'pptx':
        Byte_contents = doc.get()['Body'].read()
            
        from pptx import Presentation
        prs = Presentation(BytesIO(Byte_contents))

        texts = []
        for i, slide in enumerate(prs.slides):
            text = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = text + shape.text
            texts.append(text)
        contents = '\n'.join(texts)
        
    elif file_type == 'txt':        
        contents = doc.get()['Body'].read().decode('utf-8')

    elif file_type == 'docx':
        Byte_contents = doc.get()['Body'].read()
            
        import docx
        doc_contents =docx.Document(BytesIO(Byte_contents))

        texts = []
        for i, para in enumerate(doc_contents.paragraphs):
            if(para.text):
                texts.append(para.text)
                # print(f"{i}: {para.text}")        
        contents = '\n'.join(texts)
            
    # print('contents: ', contents)
    new_contents = str(contents).replace("\n"," ") 
    print('length: ', len(new_contents))

    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=100,
        separators=["\n\n", "\n", ".", " ", ""],
        length_function = len,
    ) 

    texts = text_splitter.split_text(new_contents) 
                
    return texts

# load csv documents from s3
def load_csv_document(s3_file_name):
    s3r = boto3.resource("s3")
    doc = s3r.Object(s3_bucket, s3_prefix+'/'+s3_file_name)

    lines = doc.get()['Body'].read().decode('utf-8').split('\n')   # read csv per line
    print('lins: ', len(lines))
        
    columns = lines[0].split(',')  # get columns
    #columns = ["Category", "Information"]  
    #columns_to_metadata = ["type","Source"]
    print('columns: ', columns)
    
    docs = []
    n = 0
    for row in csv.DictReader(lines, delimiter=',',quotechar='"'):
        # print('row: ', row)
        #to_metadata = {col: row[col] for col in columns_to_metadata if col in row}
        values = {k: row[k] for k in columns if k in row}
        content = "\n".join(f"{k.strip()}: {v.strip()}" for k, v in values.items())
        doc = Document(
            page_content=content,
            metadata={
                'name': s3_file_name,
                'page': n+1,
                'uri': path+parse.quote(s3_file_name)
            }
            #metadata=to_metadata
        )
        docs.append(doc)
        n = n+1
    print('docs[0]: ', docs[0])

    return docs

def get_summary(chat, docs):    
    text = ""
    for doc in docs:
        text = text + doc
    
    if isKorean(text)==True:
        system = (
            "다음의 <article> tag안의 문장을 요약해서 500자 이내로 설명하세오."
        )
    else: 
        system = (
            "Here is pieces of article, contained in <article> tags. Write a concise summary within 500 characters."
        )
        #prompt_template = """\n\nHuman: 다음 텍스트를 간결하게 요약하세오. 텍스트의 요점을 다루는 글머리 기호로 응답을 반환합니다.
        #prompt_template = """\n\nHuman: 아래 <text>는 문서에서 추출한 텍스트입니다. 친절한 AI Assistant로서 아래와 같이 정리해주세요.
        
        #- 50자 미안의 제목을 <title>Name: </title> 안에 넣을것
        #- 300자 미안의 설명을 <description>설명: </description> 안에 넣을것
        #- 500자 미만의 내용 요약을 <summarization>요약: </summarization> 안에 넣을것
        #- 10자 미안의 애용과 과련된 테그 5개를 <tag></tag> 테그 안에 생성할 것

        #모든 생성 결과는 한국어로 해주세요. 결과에 개행문자인 "\m"과 글자 수와 같은 부가정보는 절대 포함하지 마세요.
        #생성이 어렵거나 해당 내용을 모르는 경우 "None"로 결과를 생성하세요.
    
    human = "<article>{text}</article>"
    
    prompt = ChatPromptTemplate.from_messages([("system", system), ("human", human)])
    print('prompt: ', prompt)
    
    chain = prompt | chat    
    try: 
        result = chain.invoke(
            {
                "text": text
            }
        )
        
        summary = result.content
        print('result of summarization: ', summary)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)                    
        raise Exception ("Not able to request to LLM")
    
    return summary

    
def load_chat_history(userId, allowTime):
    dynamodb_client = boto3.client('dynamodb')

    response = dynamodb_client.query(
        TableName=callLogTableName,
        KeyConditionExpression='user_id = :userId AND request_time > :allowTime',
        ExpressionAttributeValues={
            ':userId': {'S': userId},
            ':allowTime': {'S': allowTime}
        }
    )
    # print('query result: ', response['Items'])

    for item in response['Items']:
        text = item['body']['S']
        msg = item['msg']['S']
        type = item['type']['S']

        if type == 'text':
            memory_chain.chat_memory.add_user_message(text)
            if len(msg) > MSG_LENGTH:
                memory_chain.chat_memory.add_ai_message(msg[:MSG_LENGTH])                          
            else:
                memory_chain.chat_memory.add_ai_message(msg) 
                
def getAllowTime():
    d = datetime.datetime.now() - datetime.timedelta(days = 2)
    timeStr = str(d)[0:19]
    print('allow time: ',timeStr)

    return timeStr

def isTyping(connectionId, requestId):    
    msg_proceeding = {
        'request_id': requestId,
        'msg': 'Proceeding...',
        'status': 'istyping'
    }
    #print('result: ', json.dumps(result))
    sendMessage(connectionId, msg_proceeding)

def readStreamMsg(connectionId, requestId, stream):
    msg = ""
    if stream:
        for event in stream:
            #print('event: ', event)
            msg = msg + event

            result = {
                'request_id': requestId,
                'msg': msg,
                'status': 'proceeding'
            }
            #print('result: ', json.dumps(result))
            sendMessage(connectionId, result)
    # print('msg: ', msg)
    return msg

def revise_question(connectionId, requestId, chat, query):    
    global history_length, token_counter_history    
    history_length = token_counter_history = 0
        
    if isKorean(query)==True :      
        system = (
            ""
        )  
        human = """이전 대화를 참조하여, 다음의 <question>의 뜻을 명확히 하는 새로운 질문을 한국어로 생성하세요. 새로운 질문은 원래 질문의 중요한 단어를 반드시 포함합니다. 결과는 <result> tag를 붙여주세요.
        
        <question>            
        {question}
        </question>"""
        
    else: 
        system = (
            ""
        )
        human = """Rephrase the follow up <question> to be a standalone question. Put it in <result> tags.
        <question>            
        {question}
        </question>"""
            
    prompt = ChatPromptTemplate.from_messages([("system", system), MessagesPlaceholder(variable_name="history"), ("human", human)])
    print('prompt: ', prompt)
    
    history = memory_chain.load_memory_variables({})["chat_history"]
    print('memory_chain: ', history)
                
    chain = prompt | chat    
    try: 
        result = chain.invoke(
            {
                "history": history,
                "question": query,
            }
        )
        generated_question = result.content
        
        revised_question = generated_question[generated_question.find('<result>')+8:len(generated_question)-9] # remove <result> tag                   
        print('revised_question: ', revised_question)
        
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)        
            
        sendErrorMessage(connectionId, requestId, err_msg)    
        raise Exception ("Not able to request to LLM")

    if debugMessageMode == 'true':  
        chat_history = ""
        for dialogue_turn in history:
            #print('type: ', dialogue_turn.type)
            #print('content: ', dialogue_turn.content)
            
            dialog = f"{dialogue_turn.type}: {dialogue_turn.content}\n"            
            chat_history = chat_history + dialog
                
        history_length = len(chat_history)
        print('chat_history length: ', history_length)
        
        token_counter_history = 0
        if chat_history:
            token_counter_history = chat.get_num_tokens(chat_history)
            print('token_size of history: ', token_counter_history)
            
        sendDebugMessage(connectionId, requestId, f"새로운 질문: {revised_question}\n * 대화이력({str(history_length)}자, {token_counter_history} Tokens)을 활용하였습니다.")
            
    return revised_question    
    # return revised_question.replace("\n"," ")

def isKorean(text):
    # check korean
    pattern_hangul = re.compile('[\u3131-\u3163\uac00-\ud7a3]+')
    word_kor = pattern_hangul.search(str(text))
    # print('word_kor: ', word_kor)

    if word_kor and word_kor != 'None':
        print('Korean: ', word_kor)
        return True
    else:
        print('Not Korean: ', word_kor)
        return False
    
def query_using_RAG_context(connectionId, requestId, chat, context, revised_question):    
    if isKorean(revised_question)==True:
        system = (
            """다음의 <context> tag안의 참고자료를 이용하여 상황에 맞는 구체적인 세부 정보를 충분히 제공합니다. Assistant의 이름은 서연이고, 모르는 질문을 받으면 솔직히 모른다고 말합니다.
            
            <context>
            {context}
            </context>"""
        )
    else: 
        system = (
            """Here is pieces of context, contained in <context> tags. Provide a concise answer to the question at the end. If you don't know the answer, just say that you don't know, don't try to make up an answer.
            
            <context>
            {context}
            </context>"""
        )
    
    human = "{input}"
    
    prompt = ChatPromptTemplate.from_messages([("system", system), ("human", human)])
    print('prompt: ', prompt)
                   
    chain = prompt | chat
    
    try: 
        isTyping(connectionId, requestId)  
        stream = chain.invoke(
            {
                "context": context,
                "input": revised_question,
            }
        )
        msg = readStreamMsg(connectionId, requestId, stream.content)    
        print('msg: ', msg)
        
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)        
            
        sendErrorMessage(connectionId, requestId, err_msg)    
        raise Exception ("Not able to request to LLM")

    return msg

from langchain.retrievers import AmazonKendraRetriever
kendraRetriever = AmazonKendraRetriever(
    index_id=kendraIndex, 
    top_k=top_k, 
    region_name=kendra_region,
    attribute_filter = {
        "EqualsTo": {      
            "Key": "_language_code",
            "Value": {
                "StringValue": "ko"
            }
        },
    },
)

def retrieve_from_kendra(query, top_k):
    if kendra_method == 'kendra_retriever':
        relevant_docs = retrieve_from_kendra_using_kendra_retriever(query, top_k)
    else: 
        relevant_docs = retrieve_from_kendra_using_custom_retriever(query, top_k)
    
    return relevant_docs

def retrieve_from_kendra_using_kendra_retriever(query, top_k):
    print('query: ', query)

    relevant_docs = []
    relevant_documents = kendraRetriever.get_relevant_documents(
        query=query,
        top_k=top_k,
    )
    #print('length of relevant_documents: ', len(relevant_documents))
    #print('relevant_documents: ', relevant_documents)

    rag_type = "kendra"
    api_type = "kendraRetriever"

    for i, document in enumerate(relevant_documents):
        #print('document.page_content:', document.page_content)
        #print('document.metadata:', document.metadata)
        print(f'## Document {i+1}: {document}')

        result_id = document.metadata['result_id']
        document_id = document.metadata['document_id']
        # source = document.metadata['source']
        title = document.metadata['title']
        excerpt = document.metadata['excerpt']

        uri = ""
        if "_source_uri" in document.metadata['document_attributes']:
            uri = document.metadata['document_attributes']['_source_uri']

        page = ""
        if "_excerpt_page_number" in document.metadata['document_attributes']:            
            page = document.metadata['document_attributes']['_excerpt_page_number']

        confidence = ""
        assessed_score = ""
            
        if page:
            doc_info = {
                "rag_type": rag_type,
                "api_type": api_type,
                "confidence": confidence,
                "metadata": {
                    "document_id": document_id,
                    "source": uri,
                    "title": title,
                    "excerpt": excerpt,
                    "document_attributes": {
                        "_excerpt_page_number": page
                    }
                },
                "assessed_score": assessed_score,
                "result_id": result_id
            }

        else: 
            doc_info = {
                "rag_type": rag_type,
                "api_type": api_type,
                "confidence": confidence,
                "metadata": {
                    "document_id": document_id,
                    "source": uri,
                    "title": title,
                    "excerpt": excerpt,
                },
                "assessed_score": assessed_score,
                "result_id": result_id
            }
            
        relevant_docs.append(doc_info)
    
    return relevant_docs        

def retrieve_from_kendra_using_custom_retriever(query, top_k):
    print('query: ', query)

    relevant_documents = kendraRetriever.get_relevant_documents(
        query=query,
        top_k=top_k,
    )
    print('length of relevant_documents: ', len(relevant_documents))
    print('relevant_documents: ', relevant_documents)    

    index_id = kendraIndex        
    kendra_client = boto3.client(
        service_name='kendra', 
        region_name=kendra_region,
        config = Config(
            retries=dict(
                max_attempts=10
            )
        )
    )

    try:
        resp =  kendra_client.retrieve(
            IndexId = index_id,
            QueryText = query,
            PageSize = top_k,      
            AttributeFilter = {
                "EqualsTo": {      
                    "Key": "_language_code",
                    "Value": {
                        "StringValue": "ko"
                    }
                },
            },      
        )
        # print('retrieve resp:', json.dumps(resp))
        query_id = resp["QueryId"]

        if len(resp["ResultItems"]) >= 1:
            relevant_docs = []
            retrieve_docs = []
            for query_result in resp["ResultItems"]:
                retrieve_docs.append(extract_relevant_doc_for_kendra(query_id=query_id, apiType="retrieve", query_result=query_result))
                # print('retrieve_docs: ', retrieve_docs)

            print('Looking for FAQ...')
            try:
                resp =  kendra_client.query(
                    IndexId = index_id,
                    QueryText = query,
                    PageSize = 4, # Maximum number of results returned for FAQ = 4 (default)
                    QueryResultTypeFilter = "QUESTION_ANSWER",  # 'QUESTION_ANSWER', 'ANSWER', "DOCUMENT"
                    AttributeFilter = {
                        "EqualsTo": {      
                            "Key": "_language_code",
                            "Value": {
                                "StringValue": "ko"
                            }
                        },
                    },      
                )
                print('query resp:', json.dumps(resp))
                query_id = resp["QueryId"]

                if len(resp["ResultItems"]) >= 1:
                    
                    for query_result in resp["ResultItems"]:
                        confidence = query_result["ScoreAttributes"]['ScoreConfidence']

                        #if confidence == 'VERY_HIGH' or confidence == 'HIGH' or confidence == 'MEDIUM': 
                        if confidence == 'VERY_HIGH' or confidence == 'HIGH': 
                            relevant_docs.append(extract_relevant_doc_for_kendra(query_id=query_id, apiType="query", query_result=query_result))

                            if len(relevant_docs) >= top_k:
                                break
                    # print('relevant_docs: ', relevant_docs)

                else: 
                    print('No result for FAQ')

            except Exception:
                err_msg = traceback.format_exc()
                print('error message: ', err_msg)
                raise Exception ("Not able to query from Kendra")

            for doc in retrieve_docs:                
                if len(relevant_docs) >= top_k:
                    break
                else:
                    relevant_docs.append(doc)
            
        else:  # fallback using query API
            print('No result for Retrieve API!')
            try:
                resp =  kendra_client.query(
                    IndexId = index_id,
                    QueryText = query,
                    PageSize = top_k,
                    #QueryResultTypeFilter = "DOCUMENT",  # 'QUESTION_ANSWER', 'ANSWER', "DOCUMENT"
                    AttributeFilter = {
                        "EqualsTo": {      
                            "Key": "_language_code",
                            "Value": {
                                "StringValue": "ko"
                            }
                        },
                    },      
                )
                print('query resp:', resp)
                query_id = resp["QueryId"]

                if len(resp["ResultItems"]) >= 1:
                    relevant_docs = []
                    for query_result in resp["ResultItems"]:
                        confidence = query_result["ScoreAttributes"]['ScoreConfidence']

                        if confidence == 'VERY_HIGH' or confidence == 'HIGH' or confidence == 'MEDIUM': 
                            relevant_docs.append(extract_relevant_doc_for_kendra(query_id=query_id, apiType="query", query_result=query_result))

                            if len(relevant_docs) >= top_k:
                                break
                    # print('relevant_docs: ', relevant_docs)

                else: 
                    print('No result for Query API. Finally, no relevant docs!')
                    relevant_docs = []

            except Exception:
                err_msg = traceback.format_exc()
                print('error message: ', err_msg)
                raise Exception ("Not able to query from Kendra")                

    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)        
        raise Exception ("Not able to retrieve from Kendra")     

    for i, rel_doc in enumerate(relevant_docs):
        print(f'## Document {i+1}: {json.dumps(rel_doc)}')  

    return relevant_docs

def extract_relevant_doc_for_kendra(query_id, apiType, query_result):
    rag_type = "kendra"
    if(apiType == 'retrieve'): # retrieve API
        excerpt = query_result["Content"]
        confidence = query_result["ScoreAttributes"]['ScoreConfidence']
        document_id = query_result["DocumentId"] 
        document_title = query_result["DocumentTitle"]
        
        document_uri = ""
        document_attributes = query_result["DocumentAttributes"]
        for attribute in document_attributes:
            if attribute["Key"] == "_source_uri":
                document_uri = str(attribute["Value"]["StringValue"])        
        if document_uri=="":  
            document_uri = query_result["DocumentURI"]

        doc_info = {
            "rag_type": rag_type,
            "api_type": apiType,
            "confidence": confidence,
            "metadata": {
                "document_id": document_id,
                "source": document_uri,
                "title": document_title,
                "excerpt": excerpt,
            },
            "assessed_score": "",
        }
            
    else: # query API
        query_result_type = query_result["Type"]
        confidence = query_result["ScoreAttributes"]['ScoreConfidence']
        document_id = query_result["DocumentId"] 
        document_title = ""
        if "Text" in query_result["DocumentTitle"]:
            document_title = query_result["DocumentTitle"]["Text"]
        document_uri = query_result["DocumentURI"]
        feedback_token = query_result["FeedbackToken"] 

        page = ""
        document_attributes = query_result["DocumentAttributes"]
        for attribute in document_attributes:
            if attribute["Key"] == "_excerpt_page_number":
                page = str(attribute["Value"]["LongValue"])

        if query_result_type == "QUESTION_ANSWER":
            question_text = ""
            additional_attributes = query_result["AdditionalAttributes"]
            for attribute in additional_attributes:
                if attribute["Key"] == "QuestionText":
                    question_text = str(attribute["Value"]["TextWithHighlightsValue"]["Text"])
            answer = query_result["DocumentExcerpt"]["Text"]
            excerpt = f"{question_text} {answer}"
            excerpt = excerpt.replace("\n"," ") 
        else: 
            excerpt = query_result["DocumentExcerpt"]["Text"]

        if page:
            doc_info = {
                "rag_type": rag_type,
                "api_type": apiType,
                "confidence": confidence,
                "metadata": {
                    "type": query_result_type,
                    "document_id": document_id,
                    "source": document_uri,
                    "title": document_title,
                    "excerpt": excerpt,
                    "document_attributes": {
                        "_excerpt_page_number": page
                    }
                },
                "assessed_score": "",
                "query_id": query_id,
                "feedback_token": feedback_token
            }
        else: 
            doc_info = {
                "rag_type": rag_type,
                "api_type": apiType,
                "confidence": confidence,
                "metadata": {
                    "type": query_result_type,
                    "document_id": document_id,
                    "source": document_uri,
                    "title": document_title,
                    "excerpt": excerpt,
                },
                "assessed_score": "",
                "query_id": query_id,
                "feedback_token": feedback_token
            }
    return doc_info

def priority_search(query, relevant_docs, bedrock_embedding):
    excerpts = []
    for i, doc in enumerate(relevant_docs):
        # print('doc: ', doc)
        excerpts.append(
            Document(
                page_content=doc['metadata']['excerpt'],
                metadata={
                    'name': doc['metadata']['title'],
                    'order':i,
                }
            )
        )  
    # print('excerpts: ', excerpts)

    embeddings = bedrock_embedding
    vectorstore_confidence = FAISS.from_documents(
        excerpts,  # documents
        embeddings  # embeddings
    )            
    rel_documents = vectorstore_confidence.similarity_search_with_score(
        query=query,
        k=top_k
    )

    docs = []
    for i, document in enumerate(rel_documents):
        print(f'## Document {i+1}: {document}')

        order = document[0].metadata['order']
        name = document[0].metadata['name']
        assessed_score = document[1]
        print(f"{order} {name}: {assessed_score}")

        relevant_docs[order]['assessed_score'] = int(assessed_score)

        if assessed_score < 200:
            docs.append(relevant_docs[order])    
    # print('selected docs: ', docs)

    return docs

def get_reference(docs):
    if kendra_method == 'kendra_retriever':
        reference = get_reference_using_kendra_retriever(docs)
    else:   # custom_retriever
        reference = get_reference_using_custom_retriever(docs)
        
    return reference

def get_reference_using_kendra_retriever(docs):
    reference = "\n\nFrom\n"    
    for i, doc in enumerate(docs):
        print(f'## Document {i+1}: {doc}')
                    
        page = ""
        if "document_attributes" in doc['metadata']:
            if "_excerpt_page_number" in doc['metadata']['document_attributes']:
                page = doc['metadata']['document_attributes']['_excerpt_page_number']
        uri = doc['metadata']['source']
        name = doc['metadata']['title']

        if page: 
            reference = reference + f"{i+1}. {page}page in <a href={uri} target=_blank>{name}</a>, {doc['rag_type']} ({doc['assessed_score']})\n"
        else:
            reference = reference + f"{i+1}. <a href={uri} target=_blank>{name}</a>, {doc['rag_type']} ({doc['assessed_score']})\n"
    return reference

def get_reference_using_custom_retriever(docs):
    reference = "\n\nFrom\n"    
    for i, doc in enumerate(docs):
        if doc['rag_type'] == 'kendra':
            if doc['api_type'] == 'retrieve': # Retrieve. socre of confidence is only avaialbe for English
                uri = doc['metadata']['source']
                name = doc['metadata']['title']
                reference = reference + f"{i+1}. <a href={uri} target=_blank>{name}</a>, {doc['rag_type']} ({doc['assessed_score']})\n"
            else: # Query
                confidence = doc['confidence']
                if ("type" in doc['metadata']) and (doc['metadata']['type'] == "QUESTION_ANSWER"):
                    excerpt = str(doc['metadata']['excerpt']).replace('"'," ") 
                    reference = reference + f"{i+1}. <a href=\"#\" onClick=\"alert(`{excerpt}`)\">FAQ ({confidence})</a>, {doc['rag_type']} ({doc['assessed_score']})\n"
                else:
                    uri = ""
                    if "title" in doc['metadata']:
                        #print('metadata: ', json.dumps(doc['metadata']))
                        name = doc['metadata']['title']
                        if name: 
                            uri = path+parse.quote(name)

                    page = ""
                    if "document_attributes" in doc['metadata']:
                        if "_excerpt_page_number" in doc['metadata']['document_attributes']:
                            page = doc['metadata']['document_attributes']['_excerpt_page_number']
                                                    
                    if page: 
                        reference = reference + f"{i+1}. {page}page in <a href={uri} target=_blank>{name}({confidence})</a>, {doc['rag_type']} ({doc['assessed_score']})\n"
                    else:
                        reference = reference + f"{i+1}. <a href={uri} target=_blank>{name} ({confidence})</a>, {doc['rag_type']} ({doc['assessed_score']})\n"
        elif doc['rag_type'] == 'opensearch' or doc['rag_type'] == 'faiss':
            print(f'## Document {i+1}: {doc}')
                    
            page = ""
            if "document_attributes" in doc['metadata']:
                if "_excerpt_page_number" in doc['metadata']['document_attributes']:
                    page = doc['metadata']['document_attributes']['_excerpt_page_number']
            uri = doc['metadata']['source']
            name = doc['metadata']['title']

            if page: 
                reference = reference + f"{i+1}. {page}page in <a href={uri} target=_blank>{name}</a>, {doc['rag_type']} ({doc['assessed_score']})\n"
            else:
                reference = reference + f"{i+1}. <a href={uri} target=_blank>{name}</a>, {doc['rag_type']} ({doc['assessed_score']})\n"                
    return reference
            
def retrieve_from_vectorstore(query, top_k, rag_type):
    print('query: ', query)

    relevant_docs = []
    if rag_type == 'faiss' and isReady:
        relevant_documents = vectorstore_faiss.similarity_search_with_score(
            query = query,
            k = top_k,
        )
        
        for i, document in enumerate(relevant_documents):
            print(f'## Document {i+1}: {document}')

            name = document[0].metadata['name']
            page = ""
            if "page" in document[0].metadata:
                page = document[0].metadata['page']
            uri = ""
            if "uri" in document[0].metadata:
                uri = document[0].metadata['uri']
            
            confidence = int(document[1])
            assessed_score = int(document[1])

            if page:
                doc_info = {
                    "rag_type": rag_type,
                    "confidence": confidence,
                    "metadata": {
                        "source": uri,
                        "title": name,
                        "excerpt": document[0].page_content,
                        "document_attributes": {
                            "_excerpt_page_number": page
                        }
                    },
                    "assessed_score": assessed_score,
                }

            else: 
                doc_info = {
                    "rag_type": rag_type,
                    "confidence": confidence,
                    "metadata": {
                        "source": uri,
                        "title": name,
                        "excerpt": document[0].page_content,
                    },
                    "assessed_score": assessed_score,
                }
            
            relevant_docs.append(doc_info)
            
    elif rag_type == 'opensearch':
        relevant_documents = vectorstore_opensearch.similarity_search_with_score(
            query = query,
            k = top_k,
        )

        for i, document in enumerate(relevant_documents):
            print(f'## Document {i+1}: {document}')

            name = document[0].metadata['name']
            page = ""
            if "page" in document[0].metadata:
                page = document[0].metadata['page']
            uri = ""
            if "uri" in document[0].metadata:
                uri = document[0].metadata['uri']

            excerpt = document[0].page_content
            confidence = str(document[1])
            assessed_score = str(document[1])

            if page:
                doc_info = {
                    "rag_type": rag_type,
                    "confidence": confidence,
                    "metadata": {
                        "source": uri,
                        "title": name,
                        "excerpt": excerpt,
                        "document_attributes": {
                            "_excerpt_page_number": page
                        }
                    },
                    "assessed_score": assessed_score,
                }
            else:
                doc_info = {
                    "rag_type": rag_type,
                    "confidence": confidence,
                    "metadata": {
                        "source": uri,
                        "title": name,
                        "excerpt": excerpt,
                    },
                    "assessed_score": assessed_score,
                }
            relevant_docs.append(doc_info)

    return relevant_docs

def retrieve_process_from_RAG(conn, query, top_k, rag_type):
    relevant_docs = []
    if rag_type == 'kendra':
        rel_docs = retrieve_from_kendra(query=query, top_k=top_k)      
        print('rel_docs (kendra): '+json.dumps(rel_docs))
    else:
        rel_docs = retrieve_from_vectorstore(query=query, top_k=top_k, rag_type=rag_type)
        print(f'rel_docs ({rag_type}): '+json.dumps(rel_docs))
    
    if(len(rel_docs)>=1):
        for doc in rel_docs:
            relevant_docs.append(doc)    
    
    conn.send(relevant_docs)
    conn.close()

def get_answer_using_RAG(chat, text, conv_type, connectionId, requestId, bedrock_embedding):
    reference = ""
    start_time_for_revise = time.time()

    # revise question
    revised_question = revise_question(connectionId, requestId, chat, text)    
    print('revised_question: ', revised_question)
    if debugMessageMode=='true':
        sendDebugMessage(connectionId, requestId, '[Debug]: '+revised_question)

    print('processing time for revise question: ', str(time.time() - start_time_for_revise))

    relevant_docs = []
    start_time_for_rag = time.time()
    if useParallelRAG == 'false':
        print('start the sequencial processing for multiple RAG')
        for reg in capabilities:            
            if reg == 'kendra':
                rel_docs = retrieve_from_kendra(query=revised_question, top_k=top_k)      
                print('rel_docs (kendra): '+json.dumps(rel_docs))
            else:
                rel_docs = retrieve_from_vectorstore(query=revised_question, top_k=top_k, rag_type=reg)
                print(f'rel_docs ({reg}): '+json.dumps(rel_docs))
                
            if(len(rel_docs)>=1):
                for doc in rel_docs:
                    relevant_docs.append(doc)
    else:
        print('start the parallel processing for multiple RAG')
            
        processes = []
        parent_connections = []
        for rag in capabilities:
            parent_conn, child_conn = Pipe()
            parent_connections.append(parent_conn)
            
            process = Process(target=retrieve_process_from_RAG, args=(child_conn, revised_question, top_k, rag))
            processes.append(process)

        for process in processes:
            process.start()
            
        for parent_conn in parent_connections:
            rel_docs = parent_conn.recv()

            if(len(rel_docs)>=1):
                for doc in rel_docs:
                    relevant_docs.append(doc)    

        for process in processes:
            process.join()
            
    print('processing time for RAG: ', str(time.time() - start_time_for_rag))
    #print('relevant_docs: ', relevant_docs)
        
    selected_relevant_docs = []
    if len(relevant_docs) >= 1:
        selected_relevant_docs = priority_search(revised_question, relevant_docs, bedrock_embedding)

    print('selected_relevant_docs: ', json.dumps(selected_relevant_docs))

    relevant_context = ""
    for document in selected_relevant_docs:
        relevant_context = relevant_context + document['metadata']['excerpt'] + "\n\n"
    print('relevant_context: ', relevant_context)

    # query using RAG context
    msg = query_using_RAG_context(connectionId, requestId, chat, relevant_context, revised_question)

    reference = ""
    if len(selected_relevant_docs)>=1 and enableReference=='true':
        reference = get_reference(selected_relevant_docs)
            
    return msg, reference

def create_metadata(bucket, key, meta_prefix, s3_prefix, uri, category, documentId):    
    title = key
    timestamp = int(time.time())

    metadata = {
        "Attributes": {
            "_category": category,
            "_source_uri": uri,
            "_version": str(timestamp),
        },
        "Title": title,
        "DocumentId": documentId,        
    }
    print('metadata: ', metadata)

    client = boto3.client('s3')
    try: 
        client.put_object(
            Body=json.dumps(metadata), 
            Bucket=bucket, 
            Key=meta_prefix+'/'+s3_prefix+'/'+key+'.metadata.json' 
        )
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)        
        raise Exception ("Not able to create meta file")

def getResponse(connectionId, jsonBody):
    userId  = jsonBody['user_id']
    # print('userId: ', userId)
    requestId  = jsonBody['request_id']
    # print('requestId: ', requestId)
    requestTime  = jsonBody['request_time']
    # print('requestTime: ', requestTime)
    type  = jsonBody['type']
    # print('type: ', type)
    body = jsonBody['body']
    # print('body: ', body)
    conv_type = jsonBody['conv_type']  # conversation type
    print('Conversation Type: ', conv_type)
    
    rag_type = ""
    if 'rag_type' in jsonBody:
        if jsonBody['rag_type']:
            rag_type = jsonBody['rag_type']  # RAG type
            print('rag_type: ', rag_type)

    global vectorstore_opensearch, vectorstore_faiss, enableReference
    global map_chain, memory_chain, isReady, debugMessageMode, selected_LLM

    reference = ""

    # Multi-LLM
    profile = profile_of_LLMs[selected_LLM]
    bedrock_region =  profile['bedrock_region']
    modelId = profile['model_id']
    print(f'selected_LLM: {selected_LLM}, bedrock_region: {bedrock_region}, modelId: {modelId}')
    # print('profile: ', profile)
    
    chat = get_chat(profile_of_LLMs, selected_LLM)    
    bedrock_embedding = get_embedding(profile_of_LLMs, selected_LLM)
        
    # allocate memory
    if userId in map_chain:  
        print('memory exist. reuse it!')        
        memory_chain = map_chain[userId]
        
    else: 
        print('memory does not exist. create new one!')
        memory_chain = ConversationBufferWindowMemory(memory_key="chat_history", output_key='answer', return_messages=True, k=10)
        map_chain[userId] = memory_chain
        
        allowTime = getAllowTime()
        load_chat_history(userId, allowTime)
        
    # rag sources
    if conv_type == 'qa':
        vectorstore_opensearch = OpenSearchVectorSearch(
            index_name = "rag-index-*", # all
            #index_name = 'rag-index-'+userId+'-*',
            is_aoss = False,
            ef_search = 1024, # 512(default)
            m=48,
            #engine="faiss",  # default: nmslib
            embedding_function = bedrock_embedding,
            opensearch_url=opensearch_url,
            http_auth=(opensearch_account, opensearch_passwd), # http_auth=awsauth,
        )
        print('isReady = ', isReady)

    start = int(time.time())    

    msg = ""
    if type == 'text' and body[:11] == 'list models':
        bedrock_client = boto3.client(
            service_name='bedrock',
            region_name=bedrock_region,
        )
        modelInfo = bedrock_client.list_foundation_models()    
        print('models: ', modelInfo)

        msg = f"The list of models: \n"
        lists = modelInfo['modelSummaries']
        
        for model in lists:
            msg += f"{model['modelId']}\n"
        
        msg += f"current model: {modelId}"
        print('model lists: ', msg)            
    else:             
        if type == 'text':
            text = body
            print('query: ', text)

            querySize = len(text)
            textCount = len(text.split())
            print(f"query size: {querySize}, words: {textCount}")

            if text == 'enableReference':
                enableReference = 'true'
                msg  = "Referece is enabled"
            elif text == 'disableReference':
                enableReference = 'false'
                msg  = "Reference is disabled"
            elif text == 'enableDebug':
                debugMessageMode = 'true'
                msg  = "Debug messages will be delivered to the client."
            elif text == 'disableDebug':
                debugMessageMode = 'false'
                msg  = "Debug messages will not be delivered to the client."
            elif text == 'clearMemory':
                memory_chain.clear()
                map_chain[userId] = memory_chain
                    
                print('initiate the chat memory!')
                msg  = "The chat memory was intialized in this session."
            else:          
                if conv_type == 'normal':      # normal
                    msg = general_conversation(connectionId, requestId, chat, text)  
                    
                elif conv_type == 'qa':   # question & answering
                    print(f'rag_type: {rag_type}')
                    msg, reference = get_answer_using_RAG(chat, text, conv_type, connectionId, requestId, bedrock_embedding)
                                                    
                memory_chain.chat_memory.add_user_message(text)  # append new diaglog
                memory_chain.chat_memory.add_ai_message(msg)
                
        elif type == 'document':
            isTyping(connectionId, requestId)

            object = body
            file_type = object[object.rfind('.')+1:len(object)]            
            print('file_type: ', file_type)

            if file_type == 'csv':
                docs = load_csv_document(object)
                texts = []
                for doc in docs:
                    texts.append(doc.page_content)
                print('texts: ', texts)

                msg = get_summary(chat, texts)

            elif file_type == 'pdf' or file_type == 'txt' or file_type == 'pptx' or file_type == 'docx':
                texts = load_document(file_type, object)

                docs = []
                for i in range(len(texts)):
                    docs.append(
                        Document(
                            page_content=texts[i],
                            metadata={
                                'name': object,
                                # 'page':i+1,
                                'uri': path+parse.quote(object)
                            }
                        )
                    )        
                print('docs[0]: ', docs[0])    
                print('docs size: ', len(docs))

                msg = get_summary(chat, texts)
            else:
                msg = "uploaded file: "+object
                                
            if conv_type == 'qa':
                start_time = time.time()
                category = "upload"
                documentId = "upload" + "-" + object
                if useParallelUpload == 'false':                    
                    print('upload to kendra: ', object)           
                    store_document_for_kendra(path, object, documentId)  # store the object into kendra

                    print('upload to faiss: ', object)                                                   
                    if isReady == False:   
                        embeddings = bedrock_embedding
                        vectorstore_faiss = FAISS.from_documents( # create vectorstore from a document
                            docs,  # documents
                            embeddings  # embeddings
                        )
                        isReady = True
                    else:
                        store_document_for_faiss(docs, vectorstore_faiss)

                    print('upload to opensearch: ', object)
                    store_document_for_opensearch(bedrock_embedding, docs, userId, documentId)
                    
                else:                    
                    p1 = Process(target=store_document_for_kendra, args=(path, object, documentId,))
                    p1.start(); p1.join()
                    
                    if file_type == 'pdf' or file_type == 'txt' or file_type == 'csv' or file_type == 'pptx' or file_type == 'docx':
                        # opensearch
                        p2 = Process(target=store_document_for_opensearch, args=(bedrock_embedding, docs, userId, documentId,))
                        p2.start(); p2.join()

                        # faiss
                        if isReady == False:   
                            embeddings = bedrock_embedding
                            vectorstore_faiss = FAISS.from_documents( # create vectorstore from a document
                                docs,  # documents
                                embeddings  # embeddings
                            )
                            isReady = True
                        else: 
                            vectorstore_faiss.add_documents(docs)       
                
                meta_prefix = "metadata"
                create_metadata(bucket=s3_bucket, key=object, meta_prefix=meta_prefix, s3_prefix=s3_prefix, uri=path+parse.quote(object), category=category, documentId=documentId)
                        
                print('processing time: ', str(time.time() - start_time))
                        
        elapsed_time = int(time.time()) - start
        print("total run time(sec): ", elapsed_time)        
        #print('msg+reference: ', msg+reference)

        item = {
            'user_id': {'S':userId},
            'request_id': {'S':requestId},
            'request_time': {'S':requestTime},
            'type': {'S':type},
            'body': {'S':body},
            'msg': {'S':msg+reference}
        }
        client = boto3.client('dynamodb')
        try:
            resp =  client.put_item(TableName=callLogTableName, Item=item)
        except Exception:
            err_msg = traceback.format_exc()
            print('error message: ', err_msg)
            raise Exception ("Not able to write into dynamodb")        
        #print('resp, ', resp)

    if selected_LLM >= len(profile_of_LLMs)-1:
        selected_LLM = 0
    else:
        selected_LLM = selected_LLM + 1

    return msg, reference

def lambda_handler(event, context):
    # print('event: ', event)
    
    msg = ""
    if event['requestContext']: 
        connectionId = event['requestContext']['connectionId']        
        routeKey = event['requestContext']['routeKey']
        
        if routeKey == '$connect':
            print('connected!')
        elif routeKey == '$disconnect':
            print('disconnected!')
        else:
            body = event.get("body", "")
            #print("data[0:8]: ", body[0:8])
            if body[0:8] == "__ping__":
                # print("keep alive!")                
                sendMessage(connectionId, "__pong__")
            else:
                print('connectionId: ', connectionId)
                print('routeKey: ', routeKey)
        
                jsonBody = json.loads(body)
                print('request body: ', json.dumps(jsonBody))

                requestId  = jsonBody['request_id']
                try:
                    msg, reference = getResponse(connectionId, jsonBody)

                    print('msg+reference: ', msg+reference)
                except Exception:
                    err_msg = traceback.format_exc()
                    print('err_msg: ', err_msg)

                    sendErrorMessage(connectionId, requestId, err_msg)    
                    raise Exception ("Not able to send a message")
                                    
                result = {
                    'request_id': requestId,
                    'msg': msg+reference,
                    'status': 'completed'
                }
                #print('result: ', json.dumps(result))
                sendMessage(connectionId, result)

    return {
        'statusCode': 200
    }
